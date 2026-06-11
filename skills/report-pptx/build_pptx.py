"""
report-pptx 빌더
- 구조화된 JSON 스펙을 입력받아 디자인된 7부 PPTX 리포트를 생성한다.
- 슬라이드 순서(고정): 표지 → 종목 개요 → 재무 요약 → 가격·추세 → 뉴스·심리 → 리스크 → 한 줄 종합
- 디자인: 다크 차콜 표지 + KB 옐로우(#FFBC00) 액센트, 그레이 화이트 본문,
  공통 푸터(종목·작성일·페이지), KPI 카드, 하이라이트 차트 — 증권사 리포트 톤.
- 한글 폰트 '맑은 고딕' 고정(라틴/동아시아/기호 모두 강제).
- 가드레일: 매수/매도·목표가 단정 표현 금지, 출처 없는 수치 미게재, 표 오버플로 시 행 축소.

CLI:
    python build_pptx.py --input spec.json --output reports/삼성전자.pptx
"""
from __future__ import annotations

import argparse
import json
import os
import re
import tempfile

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


# ── 디자인 토큰 ────────────────────────────────────────────────
KB_YELLOW = RGBColor(0xFF, 0xBC, 0x00)   # 포인트색
KB_YELLOW_DIM = RGBColor(0xE0, 0xA8, 0x00)
INK = RGBColor(0x21, 0x21, 0x21)         # 본문 잉크
CHARCOAL = RGBColor(0x1E, 0x1E, 0x1E)    # 표지/표 헤더 다크
SUBTLE = RGBColor(0x75, 0x75, 0x75)      # 보조 그레이
FAINT = RGBColor(0xA8, 0xA8, 0xA8)       # 푸터/캡션 그레이
HAIRLINE = RGBColor(0xDE, 0xDE, 0xDA)    # 옅은 선
PANEL = RGBColor(0xF6, 0xF6, 0xF3)       # 패널/교대 행 (그레이 화이트)
CARD = RGBColor(0xFA, 0xFA, 0xF7)        # KPI 카드 배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COVER_SUB = RGBColor(0xBD, 0xBD, 0xBD)   # 표지 보조 텍스트

KFONT = "맑은 고딕"                       # 한글 폰트 고정(글자 깨짐 방지)
MPL_FONT = "Malgun Gothic"               # matplotlib 상의 동일 폰트명

DISCLAIMER = "본 자료는 학습용 분석이며 투자 권유가 아닙니다."

# 표 한 슬라이드 데이터 행 상한(헤더 제외). 초과 시 축소 + 생략 표시.
MAX_TABLE_ROWS = 12

# 단정 표현 가드레일: 발견 시 빌드 중단. '순매수'/'매도물량' 등 정상어는 비포함.
BANNED_PATTERNS = [
    r"목표가",
    r"매수\s*추천", r"매도\s*추천",
    r"매수\s*의견", r"매도\s*의견",
    r"비중\s*확대", r"비중\s*축소",
    r"적극\s*매수", r"강력\s*매수",
    r"저점\s*매수", r"지금이\s*기회",
    r"불타기", r"물타기",
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.8)
CONTENT_W = Inches(11.73)


# ── 폰트/텍스트 헬퍼 ──────────────────────────────────────────
def _apply_font(run, size, bold, color, spacing=None):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = KFONT
    # 라틴/동아시아/기호 모두 같은 글꼴로 고정 → 한글 깨짐 방지
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", KFONT)
    if spacing is not None:
        rPr.set("spc", str(spacing))  # 1/100pt 단위 자간


def _add_text(slide, left, top, width, height, lines,
              size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, line_spacing=1.18, spacing=None):
    """lines: str 또는 (text, opts) 튜플들의 리스트. 각 항목이 한 단락."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]

    for i, item in enumerate(lines):
        text, opts = (item, {}) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        if i > 0:
            p.space_before = Pt(opts.get("space_before", 4))
        run = p.add_run()
        run.text = text
        _apply_font(run, opts.get("size", size), opts.get("bold", bold),
                    opts.get("color", color), opts.get("spacing", spacing))
    return box


def _add_rich(slide, left, top, width, height, paragraphs,
              anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """paragraphs: [{runs: [(text, size, bold, color)], space_before, align}]"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        p.line_spacing = para.get("line_spacing", line_spacing)
        if i > 0:
            p.space_before = Pt(para.get("space_before", 6))
        for text, size, bold, color in para["runs"]:
            run = p.add_run()
            run.text = text
            _apply_font(run, size, bold, color)
    return box


def _rect(slide, left, top, width, height, color, line_color=None,
          line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(0.75)
    shape.shadow.inherit = False
    return shape


def _slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ── 공통 크롬: 헤더/푸터 ──────────────────────────────────────
def _header(slide, section_no, title, subtitle=None):
    """섹션 번호 키커 + 타이틀 + 옐로우 언더라인."""
    _add_text(slide, MARGIN, Inches(0.42), Inches(3), Inches(0.3),
              f"{section_no:02d}", size=11, bold=True, color=KB_YELLOW_DIM,
              spacing=200)
    _add_text(slide, MARGIN, Inches(0.68), CONTENT_W, Inches(0.65),
              title, size=24, bold=True, color=INK)
    _rect(slide, MARGIN, Inches(1.38), Inches(0.62), Pt(2.6), KB_YELLOW)
    if subtitle:
        _add_text(slide, MARGIN + Inches(0.85), Inches(1.24), Inches(10.8),
                  Inches(0.32), subtitle, size=11, color=SUBTLE)


def _footer(slide, spec, page_no):
    _rect(slide, MARGIN, Inches(7.02), CONTENT_W, Pt(0.9), HAIRLINE)
    name = spec.get("stock_name", "")
    ticker = spec.get("ticker", "")
    label = f"{name} ({ticker})" if ticker else name
    _add_text(slide, MARGIN, Inches(7.08), Inches(7), Inches(0.32),
              f"{label} · 종목 리서치", size=8.5, color=FAINT)
    _add_text(slide, Inches(8.5), Inches(7.08), Inches(4.03), Inches(0.32),
              f"작성일 {spec.get('as_of_date', '')}   |   {page_no:02d}",
              size=8.5, color=FAINT, align=PP_ALIGN.RIGHT)


def _source_note(slide, top, text, left=None, width=None):
    _add_text(slide, left or MARGIN, top, width or CONTENT_W, Inches(0.32),
              f"출처: {text}", size=9, color=FAINT)


# ── 가드레일/정제 ─────────────────────────────────────────────
def _collect_text(spec):
    chunks = []
    chunks.extend(spec.get("overview", []))
    chunks.append(spec.get("conclusion", ""))
    price = spec.get("price", {})
    chunks.extend(price.get("comment", []))
    news = spec.get("news", {})
    chunks.extend(news.get("items", []))
    chunks.append(news.get("sentiment", ""))
    risk = spec.get("risk", {})
    chunks.extend(risk.get("items", []))
    chunks.extend(risk.get("monitoring", []))
    return "\n".join(c for c in chunks if c)


def _check_guardrails(spec):
    text = _collect_text(spec)
    hits = []
    for pat in BANNED_PATTERNS:
        for m in re.finditer(pat, text):
            hits.append(m.group(0))
    if hits:
        raise ValueError(
            "단정 표현 가드레일 위반(매수/매도·목표가 등): "
            + ", ".join(sorted(set(hits)))
            + " — 판단 근거까지만 서술하도록 문구를 수정하십시오."
        )


def _clean_metrics(metrics):
    """출처 없는 수치 지표는 게재하지 않는다."""
    return [m for m in (metrics or []) if str(m.get("source", "")).strip()]


# ── 차트 ──────────────────────────────────────────────────────
def _render_chart(chart):
    """그레이 바 + 핵심 항목(현재가)만 옐로우 하이라이트."""
    plt.rcParams["font.family"] = MPL_FONT
    plt.rcParams["axes.unicode_minus"] = False

    labels = chart["labels"]
    values = chart["values"]
    unit = chart.get("unit", "")
    horizontal = chart.get("orientation", "h") == "h"

    hi = chart.get("highlight")
    if hi is None:
        hi = next((i for i, l in enumerate(labels) if "현재" in str(l)), 0)
    colors = ["#D8D8D4"] * len(values)
    edge = ["#C8C8C4"] * len(values)
    if 0 <= hi < len(values):
        colors[hi] = "#FFBC00"
        edge[hi] = "#E0A800"

    fig, ax = plt.subplots(figsize=(6.4, 4.2))
    fig.patch.set_facecolor("white")
    if horizontal:
        y = range(len(labels))
        ax.barh(list(y), values, color=colors, edgecolor=edge, height=0.62)
        ax.set_yticks(list(y))
        ax.set_yticklabels(labels)
        ax.invert_yaxis()
        for i, v in enumerate(values):
            ax.text(v, i, f" {v:,.0f}", va="center", ha="left",
                    fontsize=9.5,
                    fontweight="bold" if i == hi else "normal",
                    color="#212121" if i == hi else "#757575")
        ax.xaxis.set_visible(False)
        ax.grid(False)
        for spine in ("top", "right", "bottom"):
            ax.spines[spine].set_visible(False)
        ax.spines["left"].set_color("#DEDEDA")
    else:
        x = range(len(labels))
        ax.bar(list(x), values, color=colors, edgecolor=edge, width=0.62)
        ax.set_xticks(list(x))
        ax.set_xticklabels(labels)
        for i, v in enumerate(values):
            ax.text(i, v, f"{v:,.0f}", va="bottom", ha="center",
                    fontsize=9.5,
                    fontweight="bold" if i == hi else "normal",
                    color="#212121" if i == hi else "#757575")
        ax.yaxis.set_visible(False)
        ax.grid(False)
        for spine in ("top", "right", "left"):
            ax.spines[spine].set_visible(False)
        ax.spines["bottom"].set_color("#DEDEDA")

    title = chart.get("title", "")
    if unit and unit not in title:
        title = f"{title}  (단위: {unit})" if title else f"단위: {unit}"
    ax.set_title(title, fontsize=12, fontweight="bold", color="#212121",
                 loc="left", pad=14)
    ax.tick_params(colors="#757575", labelsize=9.5, length=0)
    fig.tight_layout()

    path = tempfile.mktemp(suffix=".png")
    fig.savefig(path, dpi=170, facecolor="white")
    plt.close(fig)
    return path


# ── 표 ────────────────────────────────────────────────────────
def _add_table(slide, left, top, width, headers, rows):
    """다크 헤더·우측 정렬 숫자·슬림 행. 오버플로 시 축소 + 생략 행."""
    data_rows = list(rows)
    omitted = 0
    if len(data_rows) > MAX_TABLE_ROWS:
        omitted = len(data_rows) - MAX_TABLE_ROWS
        data_rows = data_rows[:MAX_TABLE_ROWS]

    n_rows = len(data_rows) + 1 + (1 if omitted else 0)
    n_cols = len(headers)
    row_h = Inches(0.4)
    height = Emu(int(row_h) * n_rows)

    gfx = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gfx.table

    # 열 폭: 첫 열(항목명) 넓게, 나머지 균등
    if n_cols > 1:
        first = Emu(int(Inches(3.9)))
        rest = Emu(int((int(width) - int(first)) / (n_cols - 1)))
        table.columns[0].width = first
        for c in range(1, n_cols):
            table.columns[c].width = rest

    for r in range(n_rows):
        table.rows[r].height = row_h

    # 헤더: 다크 차콜 + 흰 글자
    for c, htext in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CHARCOAL
        _set_cell(cell, htext, size=11.5, bold=True, color=WHITE,
                  align=PP_ALIGN.RIGHT if c else PP_ALIGN.LEFT)

    # 본문: 교대 행 + 숫자 우측 정렬
    for r, row in enumerate(data_rows, start=1):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else PANEL
            val = str(row[c]) if c < len(row) else ""
            _set_cell(cell, val, size=11, bold=(c == 0), color=INK,
                      align=PP_ALIGN.RIGHT if c else PP_ALIGN.LEFT)

    # 생략 표시
    if omitted:
        rr = n_rows - 1
        merged = table.cell(rr, 0)
        merged.merge(table.cell(rr, n_cols - 1))
        merged.fill.solid()
        merged.fill.fore_color.rgb = PANEL
        _set_cell(merged, f"… 외 {omitted}개 항목 생략(슬라이드 표시 한도)",
                  size=10, bold=False, color=SUBTLE, align=PP_ALIGN.CENTER)
    return gfx, n_rows


def _set_cell(cell, text, size, bold, color, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)
    cell.margin_left = Pt(10)
    cell.margin_right = Pt(10)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _apply_font(run, size, bold, color)


# ── 불릿(옐로우 마커 런) ─────────────────────────────────────
def _bullets(slide, left, top, width, height, items, size=14,
             color=INK, space_before=10):
    paras = []
    for it in items:
        paras.append({
            "runs": [("▪  ", size, True, KB_YELLOW_DIM),
                     (str(it), size, False, color)],
            "space_before": space_before,
        })
    if not paras:
        paras = [{"runs": [("내용 없음", size, False, SUBTLE)]}]
    return _add_rich(slide, left, top, width, height, paras)


# ── KPI 카드 ──────────────────────────────────────────────────
def _kpi_card(slide, left, top, width, height, label, value, source):
    _rect(slide, left, top, width, height, CARD,
          line_color=HAIRLINE, line_width=Pt(0.75))
    _rect(slide, left, top, width, Pt(2.4), KB_YELLOW)
    _add_rich(slide, left + Inches(0.16), top + Inches(0.1),
              width - Inches(0.32), height - Inches(0.18), [
        {"runs": [(label, 10, False, SUBTLE)]},
        {"runs": [(value, 17, True, INK)], "space_before": 2},
        {"runs": [(f"출처: {source}", 7.5, False, FAINT)], "space_before": 2},
    ])


# ── 슬라이드 ──────────────────────────────────────────────────
def _slide_cover(prs, spec):
    slide = _slide(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, CHARCOAL)          # 풀블리드 다크
    _rect(slide, 0, 0, Inches(0.16), SLIDE_H, KB_YELLOW)    # 좌측 액센트 밴드

    _add_text(slide, Inches(1.15), Inches(1.9), Inches(8), Inches(0.4),
              "EQUITY RESEARCH", size=12, bold=True, color=KB_YELLOW,
              spacing=320)

    name = spec.get("stock_name", "종목")
    _add_text(slide, Inches(1.12), Inches(2.45), Inches(11.2), Inches(1.3),
              name, size=46, bold=True, color=WHITE)

    ticker = spec.get("ticker", "")
    market = spec.get("market", "")
    sub = "  ·  ".join(x for x in [ticker, market] if x)
    if sub:
        _add_text(slide, Inches(1.15), Inches(3.6), Inches(11), Inches(0.5),
                  sub, size=16, color=COVER_SUB)

    _rect(slide, Inches(1.15), Inches(4.35), Inches(2.4), Pt(2.2), KB_YELLOW)
    _add_text(slide, Inches(1.15), Inches(4.55), Inches(11), Inches(0.45),
              "종목 리서치 리포트", size=15, color=COVER_SUB)

    _add_rich(slide, Inches(1.15), Inches(6.45), Inches(11), Inches(0.7), [
        {"runs": [("작성일  ", 11, False, COVER_SUB),
                  (spec.get("as_of_date", ""), 11, True, WHITE),
                  ("        " + DISCLAIMER, 9, False,
                   RGBColor(0x8A, 0x8A, 0x8A))]},
    ])


def _slide_overview(prs, spec):
    slide = _slide(prs)
    _header(slide, 1, "종목 개요")
    _bullets(slide, MARGIN, Inches(2.0), CONTENT_W, Inches(4.6),
             spec.get("overview", []), size=15, space_before=14)
    _footer(slide, spec, 2)


def _slide_financials(prs, spec):
    slide = _slide(prs)
    fin = spec.get("financials", {})
    _header(slide, 2, "재무 요약", "최근 3개년")
    headers = fin.get("headers", [])
    rows = fin.get("rows", [])
    if rows and not str(fin.get("source", "")).strip():
        raise ValueError("재무 요약표에는 출처가 필수입니다(출처 없는 수치 미게재).")
    if headers and rows:
        _, n_rows = _add_table(slide, MARGIN, Inches(2.0), CONTENT_W,
                               headers, rows)
        note_top = min(Inches(2.0) + Emu(int(Inches(0.4)) * n_rows)
                       + Inches(0.12), Inches(6.55))
        _source_note(slide, note_top, fin.get("source", ""))
    else:
        _add_text(slide, MARGIN, Inches(2.2), CONTENT_W, Inches(1),
                  "재무 데이터 없음", size=14, color=SUBTLE)
    _footer(slide, spec, 3)


def _slide_price(prs, spec):
    slide = _slide(prs)
    _header(slide, 3, "가격 · 추세")
    price = spec.get("price", {})
    chart = price.get("chart")
    metrics = _clean_metrics(price.get("metrics"))
    comment = price.get("comment", [])

    has_chart = bool(chart and str(chart.get("source", "")).strip()
                     and chart.get("values"))

    # KPI 카드 행 (최대 4장)
    cards = metrics[:4]
    card_bottom = Inches(2.0)
    if cards:
        gap = Inches(0.18)
        cw = Emu(int((int(CONTENT_W) - int(gap) * (len(cards) - 1)) / len(cards)))
        ch = Inches(1.08)
        for i, m in enumerate(cards):
            left = Emu(int(MARGIN) + i * (int(cw) + int(gap)))
            _kpi_card(slide, left, Inches(2.0), cw, ch,
                      m["label"], m["value"], m["source"])
        card_bottom = Inches(2.0) + Inches(1.08) + Inches(0.25)

    if has_chart:
        img = _render_chart(chart)
        slide.shapes.add_picture(img, MARGIN, card_bottom,
                                 height=Inches(3.35))
        _source_note(slide, card_bottom + Inches(3.4),
                     chart.get("source", ""), left=MARGIN, width=Inches(6))
        # 코멘트: 차트 우측
        if comment:
            cx = Inches(6.6)
            _add_text(slide, cx, card_bottom + Inches(0.1), Inches(2),
                      Inches(0.3), "추세 코멘트", size=11, bold=True,
                      color=SUBTLE, spacing=100)
            _bullets(slide, cx, card_bottom + Inches(0.5),
                     Inches(5.93), Inches(2.9), comment, size=12.5,
                     space_before=10)
    elif comment:
        _bullets(slide, MARGIN, card_bottom, CONTENT_W, Inches(3.5),
                 comment, size=13)
    # 남은 metrics(5번째 이후)는 텍스트로
    extra = metrics[4:]
    if extra and not has_chart:
        paras = [{"runs": [(f"{m['label']}: {m['value']}", 12, True, INK),
                           (f"   (출처: {m['source']})", 8.5, False, FAINT)],
                  "space_before": 6} for m in extra]
        _add_rich(slide, MARGIN, Inches(5.4), CONTENT_W, Inches(1.4), paras)
    _footer(slide, spec, 4)


def _slide_news(prs, spec):
    slide = _slide(prs)
    _header(slide, 4, "뉴스 · 심리")
    news = spec.get("news", {})
    _bullets(slide, MARGIN, Inches(1.95), CONTENT_W, Inches(3.5),
             news.get("items", []), size=12.5, space_before=12)
    sentiment = news.get("sentiment", "")
    if sentiment:
        top = Inches(5.75)
        _rect(slide, MARGIN, top, CONTENT_W, Inches(0.98), PANEL)
        _rect(slide, MARGIN, top, Inches(0.09), Inches(0.98), KB_YELLOW)
        _add_rich(slide, MARGIN + Inches(0.3), top, Inches(11.1),
                  Inches(0.98), [
            {"runs": [("시장 심리   ", 11, True, KB_YELLOW_DIM),
                      (sentiment, 12.5, True, INK)]},
        ], anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, spec, 5)


def _slide_risk(prs, spec):
    slide = _slide(prs)
    _header(slide, 5, "리스크")
    risk = spec.get("risk", {})
    items = risk.get("items", [])

    top = Inches(2.0)
    for i, item in enumerate(items[:4], 1):
        _add_rich(slide, MARGIN, top, CONTENT_W, Inches(0.62), [
            {"runs": [(f"{i:02d}", 19, True, KB_YELLOW),
                      ("   ", 19, False, INK),
                      (f"리스크 {i}.  ", 13.5, True, INK),
                      (str(item), 13.5, False, INK)]},
        ], anchor=MSO_ANCHOR.MIDDLE)
        top = top + Inches(0.72)

    mon = risk.get("monitoring", [])
    if mon:
        top = top + Inches(0.25)
        _rect(slide, MARGIN, top, Inches(0.14), Inches(0.14), KB_YELLOW)
        _add_text(slide, MARGIN + Inches(0.26), top - Inches(0.06),
                  Inches(6), Inches(0.35), "모니터링 포인트",
                  size=12.5, bold=True, color=INK)
        _bullets(slide, MARGIN, top + Inches(0.4), CONTENT_W,
                 Inches(6.8) - top, mon, size=11.5, color=SUBTLE,
                 space_before=7)
    _footer(slide, spec, 6)


def _slide_conclusion(prs, spec):
    slide = _slide(prs)
    _header(slide, 6, "종합 의견")
    top = Inches(2.55)
    h = Inches(2.1)
    _rect(slide, MARGIN, top, CONTENT_W, h, PANEL)
    _rect(slide, MARGIN, top, Inches(0.12), h, KB_YELLOW)
    _add_text(slide, MARGIN + Inches(0.45), top, Inches(10.9), h,
              spec.get("conclusion", ""), size=17, bold=True, color=INK,
              anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.35)
    _add_text(slide, MARGIN, Inches(6.55), CONTENT_W, Inches(0.4),
              DISCLAIMER, size=10.5, color=SUBTLE, align=PP_ALIGN.CENTER)
    _footer(slide, spec, 7)


# ── 엔트리포인트 ──────────────────────────────────────────────
def build(spec: dict, output_path: str) -> str:
    _check_guardrails(spec)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, spec)
    _slide_overview(prs, spec)
    _slide_financials(prs, spec)
    _slide_price(prs, spec)
    _slide_news(prs, spec)
    _slide_risk(prs, spec)
    _slide_conclusion(prs, spec)

    out_dir = os.path.dirname(output_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)
    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="report-pptx 빌더")
    parser.add_argument("--input", required=True, help="구조화 JSON 스펙 경로")
    parser.add_argument("--output", required=True, help="출력 PPTX 경로")
    args = parser.parse_args()

    with open(args.input, encoding="utf-8") as f:
        spec = json.load(f)
    path = build(spec, args.output)
    print(f"리포트 생성 완료: {path}")


if __name__ == "__main__":
    main()
